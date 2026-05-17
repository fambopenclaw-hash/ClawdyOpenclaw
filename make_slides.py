from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_PURPLE   = RGBColor(0x3B, 0x1A, 0x6B)   # header / title bg
TEAL          = RGBColor(0x00, 0x92, 0x7D)   # Petronas accent
TEAL_LIGHT    = RGBColor(0x00, 0xB0, 0xF0)   # lighter teal for sub-headers
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY      = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY     = RGBColor(0x44, 0x44, 0x44)
BLACK         = RGBColor(0x00, 0x00, 0x00)
GOLD          = RGBColor(0xF5, 0xA6, 0x23)   # highlight accent
LIGHT_PURPLE  = RGBColor(0x6A, 0x3D, 0x9F)   # gradient header text bg

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

# ── Helper: solid fill shape ─────────────────────────────────────────────────
def fill_shape(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

# ── Helper: add rectangle ────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, rgb, alpha=1.0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    fill_shape(shape, rgb)
    shape.line.fill.background()
    return shape

# ── Helper: add text box ─────────────────────────────────────────────────────
def add_tbox(slide, text, x, y, w, h,
             font_name="Calibri", font_size=18, bold=False,
             color=BLACK, align=PP_ALIGN.LEFT,
             valign=1, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline_tbox(slide, lines, x, y, w, h,
                       font_name="Calibri", font_size=16,
                       bold=False, color=BLACK,
                       align=PP_ALIGN.LEFT, line_spacing=None):
    """lines = list of (text, bold_override, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold_pt, fsize, clr) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(fsize if fsize else font_size)
        run.font.bold = bold_pt if bold_pt is not None else bold
        run.font.color.rgb = clr if clr else color
    return txBox

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════════════════
def make_title_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)

    # full dark purple bg
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_PURPLE)

    # teal accent bar left edge
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)

    # top teal stripe
    add_rect(slide, 0, 0, 13.33, 0.12, TEAL)

    # bottom footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))

    # gold thin line above footer
    add_rect(slide, 0.3, 7.0, 12.73, 0.04, GOLD)

    # Petronas "PETRONAS" bottom left
    add_tbox(slide, "PETRONAS", 0.35, 7.1, 2.5, 0.35,
             font_size=13, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

    # date bottom right
    add_tbox(slide, "© 2025 Petroliam Nasional Berhad (PETRONAS)  |  5",
             8.5, 7.12, 4.5, 0.3,
             font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    # main title
    add_tbox(slide, "REVISED PLAN",
             0.5, 1.2, 12.5, 1.1,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # sub-title
    add_tbox(slide, "@ 14th May 2026",
             0.5, 2.3, 12.5, 0.7,
             font_size=30, bold=False, color=GOLD, align=PP_ALIGN.CENTER)

    # divider
    add_rect(slide, 4.5, 3.1, 4.33, 0.05, TEAL)

    # category label box
    add_rect(slide, 4.4, 3.35, 4.53, 0.62, TEAL)
    add_tbox(slide, "SUBSURFACE",
             4.4, 3.35, 4.53, 0.62,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # subtitle tagline
    add_tbox(slide,
             "ERMAI Programme  •  LWD Technology  •  Log QC  •  BCO M1",
             0.5, 4.2, 12.5, 0.5,
             font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – OVERVIEW / AGENDA
# ════════════════════════════════════════════════════════════════════════════
def make_overview_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)

    # header bar
    add_rect(slide, 0, 0, 13.33, 1.15, DARK_PURPLE)
    add_rect(slide, 0, 0, 0.18, 1.15, TEAL)
    add_tbox(slide, "OVERVIEW", 0.4, 0.22, 6, 0.7,
             font_size=28, bold=True, color=WHITE)
    add_rect(slide, 10.5, 0.3, 2.5, 0.55, TEAL)
    add_tbox(slide, "SUBSURFACE", 10.5, 0.3, 2.5, 0.55,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))
    add_tbox(slide, "PETRONAS  |  © 2025 Petroliam Nasional Berhad (PETRONAS)",
             0.35, 7.1, 8, 0.35, font_size=9, color=MID_GRAY)

    items = [
        ("01", "ERMAI Log Conditioning – Jurassic",
         "Expedite PP deliverables; slide deck + SPE paper"),
        ("02", "LWD and ERMAI Introduction",
         "Promote LWD & ERMAI RTP to host authority for drilling & DA efficiency"),
        ("03", "ERMAI Interpretation (QC)",
         "QC YZR clastic & carbonate logs; reduce redo interpretation time"),
        ("04", "ERMAI BCO (M1)",
         "KL-led BCO milestone for M1 deployment"),
    ]
    start_y = 1.45
    for i, (num, title, desc) in enumerate(items):
        y = start_y + i * 1.45
        # number box
        add_rect(slide, 0.35, y, 0.7, 0.7, TEAL)
        add_tbox(slide, num, 0.35, y, 0.7, 0.7,
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title bar
        add_rect(slide, 1.15, y, 11.83, 0.7, DARK_PURPLE)
        add_tbox(slide, title, 1.25, y, 10, 0.7,
                 font_size=18, bold=True, color=WHITE)
        # description
        add_tbox(slide, desc, 1.25, y + 0.68, 10.5, 0.55,
                 font_size=13, color=DARK_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – STRATEGY 01
# ════════════════════════════════════════════════════════════════════════════
def make_s01_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)

    # header
    add_rect(slide, 0, 0, 13.33, 1.15, DARK_PURPLE)
    add_rect(slide, 0, 0, 0.18, 1.15, TEAL)
    add_rect(slide, 0.35, 0.22, 0.65, 0.7, TEAL)
    add_tbox(slide, "01", 0.35, 0.22, 0.65, 0.7,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "ERMAI Log Conditioning – Jurassic",
             1.1, 0.22, 9.5, 0.7, font_size=24, bold=True, color=WHITE)
    add_rect(slide, 11.3, 0.28, 1.8, 0.58, TEAL)
    add_tbox(slide, "SUBSURFACE", 11.3, 0.28, 1.8, 0.58,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # owner badge
    add_rect(slide, 0.35, 1.35, 1.2, 0.45, TEAL)
    add_tbox(slide, "Owner", 0.35, 1.35, 1.2, 0.45,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "Fahmi", 1.65, 1.38, 2, 0.4,
             font_size=16, bold=True, color=DARK_PURPLE)

    # Objective card
    add_rect(slide, 0.35, 1.95, 12.63, 0.85, WHITE)
    add_tbox(slide, "OBJECTIVE", 0.5, 2.0, 2.5, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide,
             "To expedite PP deliverables for Jurassic study",
             0.5, 2.3, 11.5, 0.45, font_size=15, bold=False, color=BLACK)

    # Key Deliverables
    add_rect(slide, 0.35, 2.95, 6.05, 0.5, DARK_PURPLE)
    add_tbox(slide, "KEY DELIVERABLES",
             0.5, 2.95, 5.8, 0.5, font_size=13, bold=True, color=WHITE)

    deliverable_items = [
        "Slide presentation — methodology & value creation",
        "SPE paper / abstract (combined with Jurassic study)",
    ]
    y = 3.55
    for item in deliverable_items:
        add_rect(slide, 0.35, y, 0.08, 0.5, TEAL)
        add_tbox(slide, item, 0.55, y, 5.8, 0.5, font_size=14, color=DARK_GRAY)
        y += 0.62

    # Notes
    add_rect(slide, 6.7, 2.95, 6.28, 0.5, DARK_PURPLE)
    add_tbox(slide, "STATUS / NOTES",
             6.85, 2.95, 5.9, 0.5, font_size=13, bold=True, color=WHITE)

    add_rect(slide, 6.7, 3.55, 6.28, 0.75, WHITE)
    add_tbox(slide, "SPE paper / abstract: IN PROGRESS",
             6.85, 3.6, 6.0, 0.65, font_size=14, bold=True,
             color=RGBColor(0xB0, 0x57, 0x00))

    # right panel — Focus
    add_rect(slide, 6.7, 4.45, 6.28, 1.4, RGBColor(0xE8, 0xF7, 0xF4))
    add_tbox(slide, "FOCUS", 6.85, 4.5, 2, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide,
             "Adoption & stakeholder buy-in\nTimely delivery of PP documents\nIntegration with wider Jurassic study",
             6.85, 4.85, 5.9, 0.9, font_size=13, color=DARK_GRAY)

    # timeline strip
    add_rect(slide, 0.35, 6.05, 12.63, 0.4, DARK_PURPLE)
    add_tbox(slide, "📅  In Progress  —  Target: End of May 2026",
             0.5, 6.05, 12.3, 0.4,
             font_size=12, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))
    add_tbox(slide, "PETRONAS  |  © 2025 Petroliam Nasional Berhad (PETRONAS)",
             0.35, 7.1, 8, 0.35, font_size=9, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – STRATEGY 02
# ════════════════════════════════════════════════════════════════════════════
def make_s02_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)

    # header
    add_rect(slide, 0, 0, 13.33, 1.15, DARK_PURPLE)
    add_rect(slide, 0, 0, 0.18, 1.15, TEAL)
    add_rect(slide, 0.35, 0.22, 0.65, 0.7, TEAL)
    add_tbox(slide, "02", 0.35, 0.22, 0.65, 0.7,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "LWD and ERMAI Introduction",
             1.1, 0.22, 9.5, 0.7, font_size=24, bold=True, color=WHITE)
    add_rect(slide, 11.3, 0.28, 1.8, 0.58, TEAL)
    add_tbox(slide, "SUBSURFACE", 11.3, 0.28, 1.8, 0.58,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # owner
    add_rect(slide, 0.35, 1.35, 1.2, 0.45, TEAL)
    add_tbox(slide, "Owner", 0.35, 1.35, 1.2, 0.45,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "Ismail / Fahmi", 1.65, 1.38, 3, 0.4,
             font_size=16, bold=True, color=DARK_PURPLE)

    # objective
    add_rect(slide, 0.35, 1.95, 12.63, 0.85, WHITE)
    add_tbox(slide, "OBJECTIVE", 0.5, 2.0, 2.5, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide,
             "Promote application of current technology (LWD & ERMAI RTP) to host authority "
             "to improve drilling & Data Acquisition efficiency",
             0.5, 2.3, 12.0, 0.45, font_size=14, color=BLACK)

    # Key Deliverables box (left column)
    add_rect(slide, 0.35, 2.95, 6.05, 0.5, DARK_PURPLE)
    add_tbox(slide, "KEY DELIVERABLES",
             0.5, 2.95, 5.8, 0.5, font_size=13, bold=True, color=WHITE)
    deliverables = [
        "LWD vs Wireline comparison pack",
        "Pilot test proposal for LWD",
        "ERMAI RTP introduction material",
        "Workshop / meeting session with TOC",
    ]
    y = 3.55
    for item in deliverables:
        add_rect(slide, 0.35, y, 0.08, 0.5, TEAL)
        add_tbox(slide, item, 0.55, y, 5.8, 0.5, font_size=13, color=DARK_GRAY)
        y += 0.58

    # Plan / Activities (right column)
    add_rect(slide, 6.7, 2.95, 6.28, 0.5, DARK_PURPLE)
    add_tbox(slide, "PLAN / ACTIVITIES",
             6.85, 2.95, 5.9, 0.5, font_size=13, bold=True, color=WHITE)

    plans = [
        "Discuss LWD application for Garraf with well team",
        "Prepare detailed LWD vs Wireline comparison (pros, time & cost savings)",
        "Convince for pilot test LWD",
        "Introduce ERMAI: LWD integration with ERMAI RTP (future autonomous operation)",
        "Conduct workshop / meeting with TOC",
    ]
    y = 3.55
    for item in plans:
        add_rect(slide, 6.7, y, 0.08, 0.5, TEAL)
        add_tbox(slide, item, 6.9, y, 5.9, 0.5, font_size=12, color=DARK_GRAY)
        y += 0.58

    # Focus badge
    add_rect(slide, 0.35, 6.05, 12.63, 0.4, TEAL)
    add_tbox(slide, "FOCUS:  Adoption & Stakeholder Buy-in  —  Target: Q2 2026",
             0.5, 6.05, 12.3, 0.4,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))
    add_tbox(slide, "PETRONAS  |  © 2025 Petroliam Nasional Berhad (PETRONAS)",
             0.35, 7.1, 8, 0.35, font_size=9, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – STRATEGY 03
# ════════════════════════════════════════════════════════════════════════════
def make_s03_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)

    # header
    add_rect(slide, 0, 0, 13.33, 1.15, DARK_PURPLE)
    add_rect(slide, 0, 0, 0.18, 1.15, TEAL)
    add_rect(slide, 0.35, 0.22, 0.65, 0.7, TEAL)
    add_tbox(slide, "03", 0.35, 0.22, 0.65, 0.7,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "ERMAI Interpretation (QC)",
             1.1, 0.22, 9.5, 0.7, font_size=24, bold=True, color=WHITE)
    add_rect(slide, 11.3, 0.28, 1.8, 0.58, TEAL)
    add_tbox(slide, "SUBSURFACE", 11.3, 0.28, 1.8, 0.58,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # owner
    add_rect(slide, 0.35, 1.35, 1.2, 0.45, TEAL)
    add_tbox(slide, "Owner", 0.35, 1.35, 1.2, 0.45,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "Ismail", 1.65, 1.38, 2, 0.4,
             font_size=16, bold=True, color=DARK_PURPLE)

    # Key Deliverables (left)
    add_rect(slide, 0.35, 1.95, 6.05, 0.5, DARK_PURPLE)
    add_tbox(slide, "KEY DELIVERABLES",
             0.5, 1.95, 5.8, 0.5, font_size=13, bold=True, color=WHITE)

    deliverable_items = [
        "ERMAI log QC output for YZR (clastic & carbonate reservoirs)",
    ]
    y = 2.55
    for item in deliverable_items:
        add_rect(slide, 0.35, y, 0.08, 0.5, TEAL)
        add_tbox(slide, item, 0.55, y, 5.8, 0.5, font_size=14, color=DARK_GRAY)
        y += 0.65

    # Plan/Activities (right)
    add_rect(slide, 6.7, 1.95, 6.28, 0.5, DARK_PURPLE)
    add_tbox(slide, "PLAN / ACTIVITIES",
             6.85, 1.95, 5.9, 0.5, font_size=13, bold=True, color=WHITE)

    plans = [
        "ERMAI log QC for YZR — clastic & carbonate reservoirs",
        "Apply QC to reduce redo interpretation time",
    ]
    y = 2.55
    for item in plans:
        add_rect(slide, 6.7, y, 0.08, 0.5, TEAL)
        add_tbox(slide, item, 6.9, y, 5.9, 0.5, font_size=13, color=DARK_GRAY)
        y += 0.62

    # Emphasis callout box
    add_rect(slide, 0.35, 3.4, 12.63, 1.5, WHITE)
    add_rect(slide, 0.35, 3.4, 0.1, 1.5, GOLD)
    add_tbox(slide, "KEY EMPHASIS", 0.6, 3.45, 3, 0.4,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide,
             "Time saving — avoid rework",
             0.6, 3.85, 5.5, 0.45,
             font_size=28, bold=True, color=DARK_PURPLE)
    add_tbox(slide,
             "QC applied to YZR logs will reduce redundant reinterpretation cycles, "
             "directly improving throughput for the subsurface team.",
             0.6, 4.45, 11.8, 0.7, font_size=13, color=DARK_GRAY)

    # Impact metrics
    add_rect(slide, 0.35, 5.05, 3.8, 1.4, RGBColor(0xE8, 0xF7, 0xF4))
    add_tbox(slide, "⏱ TIME SAVING", 0.5, 5.1, 3.5, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide, "Reduced Redo Cycles", 0.5, 5.5, 3.5, 0.5,
             font_size=16, bold=True, color=DARK_PURPLE)
    add_tbox(slide, "Avoids rework on YZR interpretation", 0.5, 6.0, 3.5, 0.4,
             font_size=11, color=DARK_GRAY)

    add_rect(slide, 4.35, 5.05, 4.3, 1.4, RGBColor(0xE8, 0xF7, 0xF4))
    add_tbox(slide, "🎯 QUALITY", 4.5, 5.1, 4.0, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide, "Consistent Log QC", 4.5, 5.5, 4.0, 0.5,
             font_size=16, bold=True, color=DARK_PURPLE)
    add_tbox(slide, "Standardised YZR clastic & carbonate", 4.5, 6.0, 4.0, 0.4,
             font_size=11, color=DARK_GRAY)

    add_rect(slide, 8.85, 5.05, 4.13, 1.4, RGBColor(0xE8, 0xF7, 0xF4))
    add_tbox(slide, "📤 OUTPUT", 9.0, 5.1, 3.9, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide, "QC Report", 9.0, 5.5, 3.9, 0.5,
             font_size=16, bold=True, color=DARK_PURPLE)
    add_tbox(slide, "For both clastic & carbonate reservoirs", 9.0, 6.0, 3.9, 0.4,
             font_size=11, color=DARK_GRAY)

    # footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))
    add_tbox(slide, "PETRONAS  |  © 2025 Petroliam Nasional Berhad (PETRONAS)",
             0.35, 7.1, 8, 0.35, font_size=9, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – STRATEGY 04
# ════════════════════════════════════════════════════════════════════════════
def make_s04_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)

    # header
    add_rect(slide, 0, 0, 13.33, 1.15, DARK_PURPLE)
    add_rect(slide, 0, 0, 0.18, 1.15, TEAL)
    add_rect(slide, 0.35, 0.22, 0.65, 0.7, TEAL)
    add_tbox(slide, "04", 0.35, 0.22, 0.65, 0.7,
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "ERMAI BCO (M1)",
             1.1, 0.22, 9.5, 0.7, font_size=24, bold=True, color=WHITE)
    add_rect(slide, 11.3, 0.28, 1.8, 0.58, TEAL)
    add_tbox(slide, "SUBSURFACE", 11.3, 0.28, 1.8, 0.58,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # owner
    add_rect(slide, 0.35, 1.35, 1.2, 0.45, TEAL)
    add_tbox(slide, "Owner", 0.35, 1.35, 1.2, 0.45,
             font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tbox(slide, "KL Lead", 1.65, 1.38, 2, 0.4,
             font_size=16, bold=True, color=DARK_PURPLE)

    # BCO milestone overview card
    add_rect(slide, 0.35, 1.95, 12.63, 2.0, WHITE)
    add_rect(slide, 0.35, 1.95, 0.1, 2.0, TEAL)
    add_tbox(slide, "MILESTONE OVERVIEW", 0.6, 2.0, 4, 0.35,
             font_size=11, bold=True, color=TEAL)
    add_tbox(slide, "Business Case Outline (BCO) — M1",
             0.6, 2.4, 11.5, 0.6,
             font_size=24, bold=True, color=DARK_PURPLE)
    add_tbox(slide,
             "The ERMAI Business Case Outline for M1 marks a key decision gate, "
             "defining the value proposition and deployment roadmap for ERMAI technology "
             "as part of Petronas's digital transformation.",
             0.6, 3.05, 11.8, 0.85, font_size=14, color=DARK_GRAY)

    # Two columns: Key Deliverables | Notes
    add_rect(slide, 0.35, 4.15, 6.05, 0.5, DARK_PURPLE)
    add_tbox(slide, "KEY DELIVERABLES",
             0.5, 4.15, 5.8, 0.5, font_size=13, bold=True, color=WHITE)
    add_rect(slide, 6.7, 4.15, 6.28, 0.5, DARK_PURPLE)
    add_tbox(slide, "NOTES",
             6.85, 4.15, 5.9, 0.5, font_size=13, bold=True, color=WHITE)

    add_rect(slide, 0.35, 4.73, 6.05, 0.9, WHITE)
    add_tbox(slide, "Details to be confirmed by KL Lead\n(pending further briefing)",
             0.5, 4.8, 5.8, 0.8, font_size=14, color=DARK_GRAY)

    add_rect(slide, 6.7, 4.73, 6.28, 0.9, WHITE)
    add_tbox(slide,
             "Status: Awaiting KL Lead update\nNext steps: internal alignment",
             6.85, 4.8, 5.9, 0.8, font_size=13, color=RGBColor(0x80, 0x80, 0x80))

    # progress indicator
    add_rect(slide, 0.35, 5.8, 12.63, 0.4, RGBColor(0xE8, 0xF7, 0xF4))
    add_rect(slide, 0.35, 5.8, 3.0, 0.4, TEAL)
    add_tbox(slide, "STATUS: PENDING",
             0.5, 5.82, 2.8, 0.35,
             font_size=12, bold=True, color=WHITE)

    # footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))
    add_tbox(slide, "PETRONAS  |  © 2025 Petroliam Nasional Berhad (PETRONAS)",
             0.35, 7.1, 8, 0.35, font_size=9, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – SUMMARY
# ════════════════════════════════════════════════════════════════════════════
def make_summary_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_PURPLE)

    # top accent
    add_rect(slide, 0, 0, 13.33, 0.12, TEAL)
    add_rect(slide, 0, 0, 0.18, 7.5, TEAL)

    # title
    add_tbox(slide, "SUBSURFACE — SUMMARY",
             0.5, 0.35, 12.5, 0.7,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.0, 1.1, 3.33, 0.05, GOLD)

    # 4 strategy summary cards
    items = [
        ("01", "ERMAI Log Conditioning – Jurassic", "Fahmi",     "In Progress",  RGBColor(0x00, 0x78, 0x5A)),
        ("02", "LWD and ERMAI Introduction",      "Ismail / Fahmi", "Q2 2026 Target", RGBColor(0x00, 0x78, 0x5A)),
        ("03", "ERMAI Interpretation (QC)",        "Ismail",    "Time Saving Focus", RGBColor(0x00, 0x78, 0x5A)),
        ("04", "ERMAI BCO (M1)",                   "KL Lead",   "Pending KL Update", RGBColor(0xB0, 0x57, 0x00)),
    ]
    cols = [(0.35, 5.35), (3.55, 8.55), (6.75, 11.75), (9.95, 13.0)]
    for i, (num, title, owner, status, clr) in enumerate(items):
        x, x2 = cols[i]
        w = x2 - x - 0.1
        # card bg
        add_rect(slide, x, 1.45, w, 5.2, RGBColor(0x4A, 0x24, 0x8F))
        # num circle
        add_rect(slide, x + w/2 - 0.35, 1.6, 0.7, 0.7, TEAL)
        add_tbox(slide, num, x + w/2 - 0.35, 1.6, 0.7, 0.7,
                 font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # title
        add_tbox(slide, title, x + 0.1, 2.45, w - 0.2, 0.9,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # owner
        add_rect(slide, x + 0.15, 3.5, w - 0.3, 0.4, RGBColor(0x3A, 0x1D, 0x6E))
        add_tbox(slide, f"👤 {owner}", x + 0.15, 3.5, w - 0.3, 0.4,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
        # status
        add_rect(slide, x + 0.15, 4.0, w - 0.3, 0.5, clr)
        add_tbox(slide, status, x + 0.15, 4.0, w - 0.3, 0.5,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # footer
    add_rect(slide, 0, 7.05, 13.33, 0.45, RGBColor(0x1A, 0x0D, 0x3D))
    add_rect(slide, 0.3, 7.0, 12.73, 0.04, GOLD)
    add_tbox(slide, "PETRONAS  |  © 2025 Petroliam Nasional Berhad (PETRONAS)",
             0.35, 7.1, 8, 0.35, font_size=9, color=MID_GRAY)
    add_tbox(slide, "REVISED PLAN @ 14th May 2026",
             10.3, 7.1, 2.7, 0.3, font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

# ── Build deck ───────────────────────────────────────────────────────────────
make_title_slide(prs)
make_overview_slide(prs)
make_s01_slide(prs)
make_s02_slide(prs)
make_s03_slide(prs)
make_s04_slide(prs)
make_summary_slide(prs)

out = "/home/fahmibakeri/.openclaw/workspace/PETRONAS_Subsurface_RevisedPlan_May2026.pptx"
prs.save(out)
print(f"Saved: {out}")