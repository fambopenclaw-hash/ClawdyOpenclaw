#!/usr/bin/env python3
"""
pptx_to_html.py — Read a .pptx file and output slide content as structured JSON.
Usage: python3 pptx_to_html.py <file.pptx>
"""
import json, sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: uv pip install python-pptx --python /tmp/reportenv/bin/python3")
    sys.exit(1)

prs = Presentation(sys.argv[1])
slides_data = []

for i, slide in enumerate(prs.slides):
    slide_content = {"slide": i + 1, "shapes": []}
    for shape in slide.shapes:
        item = {}
        if hasattr(shape, "text") and shape.text.strip():
            item["text"] = shape.text.strip()
        if shape.has_table:
            table_data = []
            for row in shape.table.rows:
                table_data.append([cell.text for cell in row.cells])
            item["table"] = table_data
        if item:
            slide_content["shapes"].append(item)
    slides_data.append(slide_content)

result = {"num_slides": len(prs.slides), "slides": slides_data}
print(json.dumps(result, indent=2, default=str))