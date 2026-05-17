#!/usr/bin/env python3
"""
make_slides.py — Generate a PETRONAS Subsurface strategy deck from JSON config.
Reads strategy data from a JSON file and outputs a formatted .pptx deck.
"""

import json, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ─────────────────────────────────────────────────────────────
DARK_PURPLE   = RGBColor(0x3B, 0x1A, 0x6B)
TEAL          = RGBColor(0x00, 0x92, 0x7D)
TEAL_LIGHT    = RGBColor(0x00, 0xB0, 0xF0)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY      = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY     = RGBColor(0x44, 0x44, 0x44)
BLACK         = RGBColor(0x00, 0x00, 0x00)
GOLD          = RGBColor(0xF5, 0xA6, 0x23)
GREEN         = RGBColor(0x00, 0x78, 0x5A)
AMBER         = RGBColor(0xB0, 0x57, 0x00)
CARD_BG       = RGBColor(0xE8, 0xF7, 0xF4)
CARD_DARK     = RGBColor(0x4A, 0x24, 0x8F)
CARD_MID      = RGBColor(0x3A, 0x1D, 0x6E)
FOOTER_DARK   = RGBColor(0x1A, 0x0D, 0x3D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def fill(shape, rgb):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb; shape.line.fill.background()

def rect(slide, x, y, w, h, rgb):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h)); fill(s, rgb); return s

def tbox(slide, text, x, y, w, h, size=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap; tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = "Calibri"
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb

def footer(slide, date_str="© 2025 Petroliam Nasional Berhad (PETRONAS)", slide_num=None):
    rect(slide, 0, 7.05, 13.33, 0.45, FOOTER_DARK)
    rect(slide, 0.3, 7.0, 12.73, 0.04, GOLD)
    tbox(slide, "PETRONAS  |  " + date_str, 0.35, 7.1, 8, 0.35, size=9, color=MID_GRAY)
    if slide_num is not None:
        tbox(slide, str(slide_num), 12.0, 7.1, 1.0, 0.3, size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def header(slide, title, brand="SUBSURFACE", badge=None, badge_color=TEAL):
    rect(slide, 0, 0, 13.33, 1.15, DARK_PURPLE)
    rect(slide, 0, 0, 0.18, 1.15, TEAL)
    tbox(slide, title, 0.4, 0.22, 9, 0.7, size=28, bold=True, color=WHITE)
    badge_text = brand if badge is None else brand
    rect(slide, 10.5, 0.3, 2.5, 0.55, badge_color)
    tbox(slide, badge_text, 10.5, 0.3, 2.5, 0.55, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 1: Title ────────────────────────────────────────────────────────────
def make_title_slide(data):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, DARK_PURPLE)
    rect(slide, 0, 0, 0.18, 7.5, TEAL)
    rect(slide, 0, 0, 13.33, 0.12, TEAL)
    footer(slide, data.get("date_str", "© 2025 Petroliam Nasional Berhad (PETRONAS)"))
    tbox(slide, data.get("title", "REVISED PLAN"), 0.5, 1.2, 12.5, 1.1, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tbox(slide, data.get("subtitle", "@ 14th May 2026"), 0.5, 2.3, 12.5, 0.7, size=30, bold=False, color=GOLD, align=PP_ALIGN.CENTER)
    rect(slide, 4.5, 3.1, 4.33, 0.05, TEAL)
    rect(slide, 4.4, 3.35, 4.53, 0.62, TEAL)
    tbox(slide, data.get("category", "SUBSURFACE"), 4.4, 3.35, 4.53, 0.62, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tagline = data.get("tagline", "ERMAI Programme  •  LWD Technology  •  Log QC  •  BCO M1")
    tbox(slide, tagline, 0.5, 4.2, 12.5, 0.5, size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── Slide 2: Overview ──────────────────────────────────────────────────────────
def make_overview_slide(items):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
    header(slide, "OVERVIEW")
    start_y = 1.45
    for i, item in enumerate(items):
        y = start_y + i * 1.45
        rect(slide, 0.35, y, 0.7, 0.7, TEAL)
        tbox(slide, item["num"], 0.35, y, 0.7, 0.7, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, 1.15, y, 11.83, 0.7, DARK_PURPLE)
        tbox(slide, item["title"], 1.25, y, 10, 0.7, size=18, bold=True, color=WHITE)
        tbox(slide, item.get("desc", ""), 1.25, y + 0.68, 10.5, 0.55, size=13, color=DARK_GRAY)
    footer(slide)

# ── Shared strategy slide shell ───────────────────────────────────────────────
def strategy_shell(num, title, owner, owner_name, subtitle=None, badge="SUBSURFACE"):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
    header(slide, title, brand=badge)
    rect(slide, 0.35, 1.35, 1.2, 0.45, TEAL)
    tbox(slide, "Owner", 0.35, 1.35, 1.2, 0.45, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tbox(slide, owner_name, 1.65, 1.38, 3, 0.4, size=16, bold=True, color=DARK_PURPLE)
    return slide

def objective_card(slide, text, y=1.95):
    rect(slide, 0.35, y, 12.63, 0.85, WHITE)
    tbox(slide, "OBJECTIVE", 0.5, y + 0.05, 2.5, 0.35, size=11, bold=True, color=TEAL)
    tbox(slide, text, 0.5, y + 0.4, 11.5, 0.4, size=15, color=BLACK)

def kd_section(slide, items, x=0.35, y=2.95, w=6.05):
    rect(slide, x, y, w, 0.5, DARK_PURPLE)
    tbox(slide, "KEY DELIVERABLES", x + 0.15, y, w - 0.3, 0.5, size=13, bold=True, color=WHITE)
    y += 0.6
    for item in items:
        rect(slide, x, y, 0.08, 0.5, TEAL)
        tbox(slide, item, x + 0.2, y, w - 0.25, 0.5, size=14, color=DARK_GRAY)
        y += 0.62

def notes_section(slide, items_or_text, x=6.7, y=2.95, w=6.28):
    rect(slide, x, y, w, 0.5, DARK_PURPLE)
    tbox(slide, "STATUS / NOTES", x + 0.15, y, w - 0.3, 0.5, size=13, bold=True, color=WHITE)
    y += 0.6
    if isinstance(items_or_text, list):
        for item in items_or_text:
            rect(slide, x, y, 0.08, 0.5, TEAL)
            tbox(slide, item, x + 0.2, y, w - 0.25, 0.5, size=13, color=DARK_GRAY)
            y += 0.58
    else:
        rect(slide, x, y, w, 0.75, WHITE)
        tbox(slide, items_or_text, x + 0.15, y + 0.05, w - 0.3, 0.65, size=14, bold=True, color=AMBER)

def focus_strip(slide, text, color=DARK_PURPLE, text_color=GOLD):
    rect(slide, 0.35, 6.05, 12.63, 0.4, color)
    tbox(slide, text, 0.5, 6.05, 12.3, 0.4, size=12, bold=True, color=text_color)

# ── Slide 3: Strategy 01 ──────────────────────────────────────────────────────
def make_s01_slide(item):
    slide = strategy_shell("01", item["title"], "Owner", item["owner"])
    objective_card(slide, item.get("objective", ""))
    kd_section(slide, item.get("deliverables", []))
    notes_section(slide, item.get("notes", []))
    rect(slide, 6.7, 4.45, 6.28, 1.4, CARD_BG)
    tbox(slide, "FOCUS", 6.85, 4.5, 2, 0.35, size=11, bold=True, color=TEAL)
    tbox(slide, item.get("focus", ""), 6.85, 4.85, 5.9, 0.9, size=13, color=DARK_GRAY)
    focus_strip(slide, item.get("timeline", "📅  In Progress  —  Target: End of May 2026"))
    footer(slide)

# ── Slide 4: Strategy 02 ──────────────────────────────────────────────────────
def make_s02_slide(item):
    slide = strategy_shell("02", item["title"], "Owner", item["owner"])
    objective_card(slide, item.get("objective", ""))
    kd_section(slide, item.get("deliverables", []), w=6.05)
    rect(slide, 6.7, 2.95, 6.28, 0.5, DARK_PURPLE)
    tbox(slide, "PLAN / ACTIVITIES", 6.85, 2.95, 5.9, 0.5, size=13, bold=True, color=WHITE)
    y = 3.55
    for p in item.get("plan", []):
        rect(slide, 6.7, y, 0.08, 0.5, TEAL)
        tbox(slide, p, 6.9, y, 5.9, 0.5, size=12, color=DARK_GRAY)
        y += 0.58
    focus_strip(slide, item.get("focus", "FOCUS:  Adoption & Stakeholder Buy-in  —  Target: Q2 2026"), color=TEAL, text_color=WHITE)
    footer(slide)

# ── Slide 5: Strategy 03 ──────────────────────────────────────────────────────
def make_s03_slide(item):
    slide = strategy_shell("03", item["title"], "Owner", item["owner"])
    kd_section(slide, item.get("deliverables", []), w=6.05, y=1.95)
    rect(slide, 6.7, 1.95, 6.28, 0.5, DARK_PURPLE)
    tbox(slide, "PLAN / ACTIVITIES", 6.85, 1.95, 5.9, 0.5, size=13, bold=True, color=WHITE)
    y = 2.55
    for p in item.get("plan", []):
        rect(slide, 6.7, y, 0.08, 0.5, TEAL)
        tbox(slide, p, 6.9, y, 5.9, 0.5, size=13, color=DARK_GRAY)
        y += 0.62
    rect(slide, 0.35, 3.4, 12.63, 1.5, WHITE)
    rect(slide, 0.35, 3.4, 0.1, 1.5, GOLD)
    tbox(slide, "KEY EMPHASIS", 0.6, 3.45, 3, 0.4, size=11, bold=True, color=TEAL)
    tbox(slide, item.get("emphasis_title", "Time saving — avoid rework"), 0.6, 3.85, 5.5, 0.45, size=28, bold=True, color=DARK_PURPLE)
    tbox(slide, item.get("emphasis_body", ""), 0.6, 4.45, 11.8, 0.7, size=13, color=DARK_GRAY)
    # 3 metric cards
    metrics = item.get("metrics", [])
    for idx, m in enumerate(metrics[:3]):
        mx = [0.35, 4.35, 8.85][idx]
        rect(slide, mx, 5.05, [3.8, 4.3, 4.13][idx], 1.4, CARD_BG)
        tbox(slide, m.get("label", ""), mx + 0.15, 5.1, 3.5, 0.35, size=11, bold=True, color=TEAL)
        tbox(slide, m.get("title", ""), mx + 0.15, 5.5, 3.5, 0.5, size=16, bold=True, color=DARK_PURPLE)
        tbox(slide, m.get("desc", ""), mx + 0.15, 6.0, 3.5, 0.4, size=11, color=DARK_GRAY)
    footer(slide)

# ── Slide 6: Strategy 04 ──────────────────────────────────────────────────────
def make_s04_slide(item):
    slide = strategy_shell("04", item["title"], "Owner", item["owner"])
    rect(slide, 0.35, 1.95, 12.63, 2.0, WHITE)
    rect(slide, 0.35, 1.95, 0.1, 2.0, TEAL)
    tbox(slide, "MILESTONE OVERVIEW", 0.6, 2.0, 4, 0.35, size=11, bold=True, color=TEAL)
    tbox(slide, item.get("milestone_title", "Business Case Outline (BCO) — M1"), 0.6, 2.4, 11.5, 0.6, size=24, bold=True, color=DARK_PURPLE)
    tbox(slide, item.get("milestone_body", ""), 0.6, 3.05, 11.8, 0.85, size=14, color=DARK_GRAY)
    rect(slide, 0.35, 4.15, 6.05, 0.5, DARK_PURPLE)
    tbox(slide, "KEY DELIVERABLES", 0.5, 4.15, 5.8, 0.5, size=13, bold=True, color=WHITE)
    rect(slide, 6.7, 4.15, 6.28, 0.5, DARK_PURPLE)
    tbox(slide, "NOTES", 6.85, 4.15, 5.9, 0.5, size=13, bold=True, color=WHITE)
    rect(slide, 0.35, 4.73, 6.05, 0.9, WHITE)
    tbox(slide, item.get("deliverables_text", "Details to be confirmed by KL Lead\n(pending further briefing)"), 0.5, 4.8, 5.8, 0.8, size=14, color=DARK_GRAY)
    rect(slide, 6.7, 4.73, 6.28, 0.9, WHITE)
    tbox(slide, item.get("notes_text", "Status: Awaiting KL Lead update\nNext steps: internal alignment"), 6.85, 4.8, 5.9, 0.8, size=13, color=MID_GRAY)
    rect(slide, 0.35, 5.8, 12.63, 0.4, CARD_BG)
    rect(slide, 0.35, 5.8, 3.0, 0.4, TEAL)
    tbox(slide, "STATUS: " + item.get("status", "PENDING"), 0.5, 5.82, 2.8, 0.35, size=12, bold=True, color=WHITE)
    footer(slide)

# ── Slide 7: Summary ───────────────────────────────────────────────────────────
def make_summary_slide(items, date_str=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, DARK_PURPLE)
    rect(slide, 0, 0, 13.33, 0.12, TEAL)
    rect(slide, 0, 0, 0.18, 7.5, TEAL)
    tbox(slide, "SUBSURFACE — SUMMARY", 0.5, 0.35, 12.5, 0.7, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, 5.0, 1.1, 3.33, 0.05, GOLD)
    cols = [(0.35, 5.35), (3.55, 8.55), (6.75, 11.75), (9.95, 13.0)]
    for i, item in enumerate(items):
        x, x2 = cols[i]; w = x2 - x - 0.1
        rect(slide, x, 1.45, w, 5.2, CARD_DARK)
        rect(slide, x + w/2 - 0.35, 1.6, 0.7, 0.7, TEAL)
        tbox(slide, item["num"], x + w/2 - 0.35, 1.6, 0.7, 0.7, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tbox(slide, item["title"], x + 0.1, 2.45, w - 0.2, 0.9, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, x + 0.15, 3.5, w - 0.3, 0.4, CARD_MID)
        tbox(slide, f"👤 {item['owner']}", x + 0.15, 3.5, w - 0.3, 0.4, size=11, color=WHITE, align=PP_ALIGN.CENTER)
        status_color = GREEN if "progress" in item.get("status","").lower() else AMBER
        rect(slide, x + 0.15, 4.0, w - 0.3, 0.5, status_color)
        tbox(slide, item.get("status", ""), x + 0.15, 4.0, w - 0.3, 0.5, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    footer(slide, date_str if date_str else "REVISED PLAN @ 14th May 2026")
    tbox(slide, "REVISED PLAN @ 14th May 2026", 10.3, 7.1, 2.7, 0.3, size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

# ── Main ────────────────────────────────────────────────────────────────────────
def main():
    config_path = Path(sys.argv[1] if len(sys.argv) > 1 else "deck_config.json")
    if not config_path.exists():
        print(f"ERROR: config file not found: {config_path}"); sys.exit(1)
    with open(config_path) as f:
        data = json.load(f)

    make_title_slide(data.get("title_slide", {}))
    make_overview_slide(data.get("strategies", []))
    for item in data.get("strategies", []):
        num = item.get("num", "01")
        if num == "01": make_s01_slide(item)
        elif num == "02": make_s02_slide(item)
        elif num == "03": make_s03_slide(item)
        elif num == "04": make_s04_slide(item)
    make_summary_slide(data.get("strategies", []), data.get("title_slide", {}).get("subtitle", ""))

    out = Path(data.get("output", "PETRONAS_Subsurface_Deck.pptx"))
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    main()